import streamlit as st
import os
import numpy as np  
import fitz  # PyMuPDF
import easyocr  # EasyOCR
import cv2
import shutil
import re
import threading
from PIL import Image
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from transformers import pipeline


# Setup EasyOCR reader
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"  # To address OpenMP runtime issue
reader = easyocr.Reader(['en'])  # Initialize EasyOCR

# Function to convert PDF pages to images
def convert_pdf_pages_to_images(pdf_path, image_folder_path):
    doc = fitz.open(pdf_path)
    if not os.path.exists(image_folder_path):
        os.makedirs(image_folder_path)
    images = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        # Increase the resolution by specifying a higher zoom factor
        # Default DPI in PDFs is usually 72, so zooming by 3 gives you 216 DPI
        zoom_x = 3.5  # horizontal zoom
        zoom_y = 3.5  # vertical zoom
        mat = fitz.Matrix(zoom_x, zoom_y)  # Zoom factor 3 in each dimension
        pix = page.get_pixmap(matrix=mat, alpha=False)  # Render page to an image
        image_path = os.path.join(image_folder_path, f"page_{page_num}.png")
        pix.save(image_path)
        images.append(image_path)
    return images

# Function to detect highlighted regions in images
def detect_highlighted_regions(image_paths):
    lower_yellow = np.array([20, 100, 100])
    upper_yellow = np.array([30, 255, 255])
    highlighted_regions = []
    for image_path in image_paths:
        image = cv2.imread(image_path)
        hsv = cv2.cvtColor(image, cv2.COLOR_BGR2HSV)
        mask = cv2.inRange(hsv, lower_yellow, upper_yellow)
        contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        contours = refine_contours(contours)  # Refine contours to focus on highlighted regions
        contours = sort_contours(contours)  # Sort contours to ensure correct order of text extraction
        for contour in contours:
            x, y, w, h = cv2.boundingRect(contour)
            highlighted_regions.append((image_path, (x, y, x+w, y+h)))
    return highlighted_regions

# Function to refine contours based on certain criteria
def refine_contours(contours):
    refined = []
    for contour in contours:
        _, _, w, h = cv2.boundingRect(contour)
        if w > 10 and h > 10:  # Example criteria
            refined.append(contour)
    return refined

# Function to sort contours from top to bottom
def sort_contours(contours):
    return sorted(contours, key=lambda c: cv2.boundingRect(c)[1])

# Function to extract text from highlighted regions
def extract_text_from_highlights(highlighted_regions):
    extracted_texts = []
    for image_path, (x1, y1, x2, y2) in highlighted_regions:
        image = Image.open(image_path).convert('RGB')
        cropped_image = image.crop((x1, y1, x2, y2))
        result = reader.readtext(np.array(cropped_image), detail=0)
        extracted_text = " ".join(result)
        extracted_texts.append(extracted_text)
    return extracted_texts


# Function to count words in a text
def count_words(text):
    words = re.findall(r'\w+', text)
    return len(words)

# Function to delete images in a folder
def delete_images(folder_path):
    files = os.listdir(folder_path)
    for file in files:
        if file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
            file_path = os.path.join(folder_path, file)
            try:
                os.remove(file_path)
                # st.write(f"Deleted: {file_path}")
            except Exception as e:
                st.write(f"Error deleting {file_path}: {e}")

# Function to save extracted text to a Word document
def save_text_to_word(extracted_texts, output_path):
    doc = Document()
    for text in extracted_texts:
        num_words = count_words(text)
        max_words = 200
        min_words = 20
        if num_words <= min_words:
            summarized_text = text
        else:
            max_length = min(max_words, num_words)
            summarizer = pipeline("summarization", model="t5-small", tokenizer="t5-small")
            summarized_text = summarizer(text, max_length=max_length, min_length=min_words)[0]['summary_text']
        doc.add_paragraph(summarized_text)
    doc.save(output_path)
    folder_path = "images"
    delete_images(folder_path)

# Function to create a presentation from a Word document
def create_presentation_from_word(doc_path, pptx_path, template_path):
    doc = Document(doc_path)
    prs = Presentation(template_path)
    single_line_heading_pattern = re.compile(r'^\d+(\.\d+)*(\.\d+)*\s+[A-Z].*')
    multi_line_heading_number_pattern = re.compile(r'^\d+(\.\d+)*$')
    word_limit_per_slide = 110

    def add_initial_content_slide():
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        _, content_shape = slide.shapes.title, slide.placeholders[1]
        customize_content_shape(content_shape)
        return content_shape.text_frame

    content_started = False
    text_frame_for_initial_content = None

    for i, paragraph in enumerate(doc.paragraphs):
        text = paragraph.text.strip()

        if not content_started and (single_line_heading_pattern.match(text) is None and not text.isdigit()):
            parts = split_text_into_parts(text, word_limit_per_slide)

            for part in parts:
                if text_frame_for_initial_content is None or parts.index(part) > 0:
                    text_frame_for_initial_content = add_initial_content_slide()
                add_content_to_slide(text_frame_for_initial_content, part)

            continue

        content_started = True
        is_multi_line_heading = (multi_line_heading_number_pattern.match(text) and
                                 i + 1 < len(doc.paragraphs) and
                                 doc.paragraphs[i + 1].text.strip()[0].islower())

        if single_line_heading_pattern.match(text) or is_multi_line_heading:
            slide_layout = prs.slide_layouts[1]
            current_slide = prs.slides.add_slide(slide_layout)
            title_shape, content_shape = current_slide.shapes.title, current_slide.placeholders[1]

            title_text = text if not is_multi_line_heading else f"{text} {doc.paragraphs[i + 1].text.strip()}"
            if is_multi_line_heading:
                i += 1

            customize_title_shape(title_shape, title_text)
            customize_content_shape(content_shape)

        elif 'current_slide' in locals():
            parts = split_text_into_parts(text.replace('_', '.'), word_limit_per_slide)
            for part in parts:
                if parts.index(part) > 0:
                    current_slide = prs.slides.add_slide(slide_layout)
                    _, content_shape = current_slide.shapes.title, current_slide.placeholders[1]
                    customize_content_shape(content_shape)
                add_content_to_slide(content_shape.text_frame, part)

    prs.save(pptx_path)

# Function to format headings in a Word document
def format_headings_in_word(doc_path, output_path):
    doc = Document(doc_path)
    single_line_heading_pattern = re.compile(r'^\d+(\.\d+)*(\.\d+)*\s+[A-Z].*')
    heading_number_pattern = re.compile(r'^\d+(\.\d+)*$')
    figure_line_pattern = re.compile(r'^Figure\s+')

    new_doc = Document()

    i = 0
    while i < len(doc.paragraphs):
        paragraph = doc.paragraphs[i]
        text = paragraph.text.strip().replace("_", ".")
        
        if figure_line_pattern.match(text):
            i += 1
            continue

        if single_line_heading_pattern.match(text) or heading_number_pattern.match(text):
            run = new_doc.add_paragraph().add_run(text)
            run.bold = True
            run.font.size = Pt(12)
            if heading_number_pattern.match(text) and i + 1 < len(doc.paragraphs) and not figure_line_pattern.match(doc.paragraphs[i + 1].text.strip()):
                i += 1
                next_text = doc.paragraphs[i].text.strip().replace('_', '.')
                run = new_doc.add_paragraph().add_run(next_text)
                run.bold = True
                run.font.size = Pt(12)
        else:
            text_content = [text]
            while i + 1 < len(doc.paragraphs) and not single_line_heading_pattern.match(doc.paragraphs[i + 1].text.strip()) and not heading_number_pattern.match(doc.paragraphs[i + 1].text.strip()) and not figure_line_pattern.match(doc.paragraphs[i + 1].text.strip()):
                i += 1
                next_text = doc.paragraphs[i].text.strip().replace('_', '.')
                text_content.append(next_text)
            consolidated_text = ' '.join(text_content)
            new_paragraph = new_doc.add_paragraph(consolidated_text)
            new_paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        i += 1

    new_doc.save(output_path)

# Function to split text into parts ensuring each part has around 'limit' words without cutting sentences in the middle
def split_text_into_parts(text, limit):
    sentences = re.split(r'(?<=[.!?])\s+', text)
    parts = []
    part_words = []
    current_count = 0

    for sentence in sentences:
        sentence_words = sentence.split()
        sentence_length = len(sentence_words)

        if current_count + sentence_length > limit:
            if part_words:
                parts.append(' '.join(part_words))
            part_words = sentence_words
            current_count = sentence_length
        else:
            part_words.extend(sentence_words)
            current_count += sentence_length

    if part_words:
        parts.append(' '.join(part_words))

    return parts

# Function to customize title shape in PowerPoint
def customize_title_shape(title_shape, title_text):
    sentences = re.split(r'(?<=\.)\s+', title_text)
    for index, sentence in enumerate(sentences):
        if sentence:
            p = title_shape.text_frame.add_paragraph() if index > 0 or title_shape.text_frame.paragraphs[0].text else title_shape.text_frame.paragraphs[0]
            p.text = sentence.replace('_', '.')
            font_size = 28
            if 41 > len(sentence) > 31:
                font_size = 24
            elif 51 > len(sentence) >= 41:
                font_size = 22
            elif len(sentence) >= 51:
                font_size = 18
            p.alignment = PP_ALIGN.JUSTIFY
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.name = 'Calibri'
    title_shape.width = Cm(21)
    title_shape.height = Cm(2.5)
    title_shape.left = Cm(0.3)
    title_shape.top = Cm(0.4)

# Function to customize content shape in PowerPoint
def customize_content_shape(content_shape):
    content_shape.width = Cm(24)
    content_shape.height = Cm(15)
    content_shape.left = Cm(0.3)
    content_shape.top = Cm(2.8)

# Function to add content to slide in PowerPoint
def add_content_to_slide(text_frame, text):
    sentences = re.split(r'(?<=\.)\s+', text)
    for index, sentence in enumerate(sentences):
        if sentence:
            p = text_frame.add_paragraph() if index > 0 or text_frame.paragraphs[0].text else text_frame.paragraphs[0]
            p.text = sentence.replace('_', '.')
            p.alignment = PP_ALIGN.JUSTIFY
            for run in p.runs:
                run.font.size = Pt(21)
                run.font.name = 'Calibri'

# Streamlit app
st.title("PDF to PPT Converter")

pdf_path = st.file_uploader("Select PDF:", type=["pdf"])
pptx_template_path = st.file_uploader("Select PowerPoint Template:", type=["pptx"])
output_dir = st.text_input("Enter Output Directory Path:")
output_dir = os.path.abspath(output_dir) if output_dir else None

if st.button("Convert PDF to PPT"):
    if pdf_path is None:
        st.warning("Please select a PDF file.")
    elif pptx_template_path is None:
        st.warning("Please select a PowerPoint template.")
    elif not output_dir:
        st.warning("Please enter the output directory path.")
    else:
        image_folder_path = os.path.join(output_dir, "images")
        pdf_filename = os.path.basename(pdf_path.name)
        pptx_template_filename = os.path.basename(pptx_template_path.name)
        output_word_file = os.path.join(output_dir, f"{os.path.splitext(pdf_filename)[0]}_extracted_text.docx")
        formatted_output_word_file = os.path.join(output_dir, f"{os.path.splitext(pdf_filename)[0]}_formatted_extracted_text.docx")
        pptx_output_path = os.path.join(output_dir, f"{os.path.splitext(pdf_filename)[0]}_presentation.pptx")
        
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Process PDF to extract text and convert to PowerPoint presentation
        try:
            st.write("Converting PDF to Word...")
            image_paths = convert_pdf_pages_to_images(pdf_path, image_folder_path)
            highlighted_regions = detect_highlighted_regions(image_paths)
            extracted_texts = extract_text_from_highlights(highlighted_regions)
            save_text_to_word(extracted_texts, output_word_file)
            
            st.write("Formatting Word document...")
            format_headings_in_word(output_word_file, formatted_output_word_file)
            
            st.write("Creating PowerPoint presentation...")
            create_presentation_from_word(formatted_output_word_file, pptx_output_path, pptx_template_path)
            delete_images(image_folder_path)
            shutil.rmtree(image_folder_path)
        except Exception as e:
            st.error(f"An error occurred: {e}")
