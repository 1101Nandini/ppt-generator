from PIL import Image
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH 
from tkinter import filedialog, messagebox
from tkinter import messagebox
from transformers import pipeline
import numpy as np
import cv2
import threading
import os
import fitz  # PyMuPDF
import easyocr  # EasyOCR
import docx
import re
import tkinter as tk
import threading

# Setup EasyOCR reader
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"  # To address OpenMP runtime issue
reader = easyocr.Reader(['en'])  # Initialize EasyOCR

# Implementation of GUI actions
def select_file(entry_widget):
    file_path = filedialog.askopenfilename(filetypes=[("PDF files", "*.pdf"), ("PowerPoint templates", "*.pptx")])
    if file_path:
        entry_widget.delete(0, tk.END)
        entry_widget.insert(0, file_path)


def select_folder(title, entry_widget):
    folder_path = filedialog.askdirectory(title=title)
    if folder_path:
        entry_widget.delete(0, tk.END)
        entry_widget.insert(0, folder_path)


def start_process_thread():
    pdf_path = entry_pdf.get()
    image_folder_path = entry_image_folder.get()
    output_dir = entry_output_dir.get()
    template_path = entry_template.get()
    threading.Thread(target=lambda: run_process(pdf_path, image_folder_path, output_dir, template_path, button_process, status_label)).start()


def convert_pdf_pages_to_images(pdf_path, image_folder_path):
    doc = fitz.open(pdf_path)
    if not os.path.exists(image_folder_path):
        os.makedirs(image_folder_path)
    images = []
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        # Increase the resolution by specifying a higher zoom factor
        # Default DPI in PDFs is usually 72, so zooming by 3 gives you 216 DPI
        zoom_x = 3.5  # horizontal zoom
        zoom_y = 3.5  # vertical zoom
        mat = fitz.Matrix(zoom_x, zoom_y)  # Zoom factor 3 in each dimension
        pix = page.get_pixmap(matrix=mat, alpha=False)  # Render page to an image
        image_path = os.path.join(image_folder_path, f"page_{page_num}.png")
        pix.save(image_path)
        images.append(image_path)
    return images


def detect_highlighted_regions(image_paths):
    lower_yellow = np.array([20, 100, 100])
    upper_yellow = np.array([30, 255, 255])
    highlighted_regions = []
    for image_path in image_paths:
        image = cv2.imread(image_path)
        hsv = cv2.cvtColor(image, cv2.COLOR_BGR2HSV)
        mask = cv2.inRange(hsv, lower_yellow, upper_yellow)
        contours, _ = cv2.findContours(mask, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
        contours = refine_contours(contours)  # Refine contours to focus on highlighted regions
        contours = sort_contours(contours)  # Sort contours to ensure correct order of text extraction
        for contour in contours:
            x, y, w, h = cv2.boundingRect(contour)
            highlighted_regions.append((image_path, (x, y, x+w, y+h)))
    return highlighted_regions


def refine_contours(contours):
    """Refine contours based on certain criteria."""
    refined = []
    for contour in contours:
        _, _, w, h = cv2.boundingRect(contour)
        if w > 10 and h > 10:  # Example criteria
            refined.append(contour)
    return refined


def sort_contours(contours):
    """Sort contours from top to bottom."""
    return sorted(contours, key=lambda c: cv2.boundingRect(c)[1])


def extract_text_from_highlights(highlighted_regions):
    extracted_texts = []
    for image_path, (x1, y1, x2, y2) in highlighted_regions:
        image = Image.open(image_path).convert('RGB')
        cropped_image = image.crop((x1, y1, x2, y2))
        result = reader.readtext(np.array(cropped_image), detail=0)
        extracted_text = " ".join(result)
        
        extracted_texts.append(extracted_text)
    return extracted_texts


# Function to delete images in a folder
def delete_images(folder_path):
    # List all files in the folder
    files = os.listdir(folder_path)
    
    # Iterate over each file
    for file in files:
        # Check if the file is an image (you can add more image file extensions as needed)
        if file.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
            # Construct the full path to the file
            file_path = os.path.join(folder_path, file)
            try:
                # Delete the file
                os.remove(file_path)
                print(f"Deleted: {file_path}")
            except Exception as e:
                print(f"Error deleting {file_path}: {e}")

# Function to count words in a text
def count_words(text):
    words = re.findall(r'\w+', text)
    return len(words)


def save_text_to_word(extracted_texts, output_path):
    doc = Document()
    for text in extracted_texts:
        num_words = count_words(text)
        max_words = 200
        min_words = 20

        if num_words <= min_words:
            summarized_text = text  # Return the original text if it has fewer words than the minimum
        else:
            max_length = min(max_words, num_words)  # Adjust the max_length based on the number of words
            summarizer = pipeline("summarization", model="t5-small", tokenizer="t5-small")
            summarized_text = summarizer(text, max_length=max_length, min_length=min_words)[0]['summary_text']
        doc.add_paragraph(summarized_text)
    doc.save(output_path)
    folder_path = "images"
    delete_images(folder_path)


def create_presentation_from_word(doc_path, pptx_path, template_path):
    doc = Document(doc_path)
    prs = Presentation(template_path)

    single_line_heading_pattern = re.compile(r'^\d+(\.\d+)*(\.\d+)*\s+[A-Z].*')
    multi_line_heading_number_pattern = re.compile(r'^\d+(\.\d+)*$')
    word_limit_per_slide = 110  # Limit to 110 words per slide

    def add_initial_content_slide():
        slide_layout = prs.slide_layouts[1]  # Layout for content slides
        slide = prs.slides.add_slide(slide_layout)
        _, content_shape = slide.shapes.title, slide.placeholders[1]
        customize_content_shape(content_shape)
        return content_shape.text_frame

    content_started = False
    text_frame_for_initial_content = None

    for i, paragraph in enumerate(doc.paragraphs):
        text = paragraph.text.strip()

        if not content_started and (single_line_heading_pattern.match(text) is None and not text.isdigit()):
            # Split the paragraph text into parts according to the word limit
            parts = split_text_into_parts(text, word_limit_per_slide)

            for part in parts:
                # If there's no text frame for initial content, or if we need a new slide for additional parts
                if text_frame_for_initial_content is None or parts.index(part) > 0:
                    # Add a new slide and set its text frame as the current one
                    text_frame_for_initial_content = add_initial_content_slide()
                    
                # Add the current part to the slide
                add_content_to_slide(text_frame_for_initial_content, part)

            continue

        content_started = True
        is_multi_line_heading = (multi_line_heading_number_pattern.match(text) and
                                 i + 1 < len(doc.paragraphs) and
                                 doc.paragraphs[i + 1].text.strip()[0].islower())

        if single_line_heading_pattern.match(text) or is_multi_line_heading:
            slide_layout = prs.slide_layouts[1]  # Layout for content slides
            current_slide = prs.slides.add_slide(slide_layout)
            title_shape, content_shape = current_slide.shapes.title, current_slide.placeholders[1]

            title_text = text if not is_multi_line_heading else f"{text} {doc.paragraphs[i + 1].text.strip()}"
            if is_multi_line_heading:
                i += 1  # Skip next paragraph as it's part of the title

            customize_title_shape(title_shape, title_text)
            customize_content_shape(content_shape)

        elif 'current_slide' in locals():
            parts = split_text_into_parts(text.replace('_', '.'), word_limit_per_slide)
            for part in parts:
                if parts.index(part) > 0:  # Create new slide for overflow content
                    current_slide = prs.slides.add_slide(slide_layout)
                    _, content_shape = current_slide.shapes.title, current_slide.placeholders[1]
                    customize_content_shape(content_shape)
                add_content_to_slide(content_shape.text_frame, part)

    prs.save(pptx_path)


def format_headings_in_word(doc_path, output_path):
    doc = Document(doc_path)
    single_line_heading_pattern = re.compile(r'^\d+(\.\d+)*(\.\d+)*\s+[A-Z].*')
    heading_number_pattern = re.compile(r'^\d+(\.\d+)*$')
    figure_line_pattern = re.compile(r'^Figure\s+')

    new_doc = Document()

    i = 0
    while i < len(doc.paragraphs):
        paragraph = doc.paragraphs[i]
        text = paragraph.text.strip().replace("_", ".")
        
        if figure_line_pattern.match(text):
            i += 1
            continue

        if single_line_heading_pattern.match(text) or heading_number_pattern.match(text):
            run = new_doc.add_paragraph().add_run(text)
            run.bold = True
            run.font.size = Pt(12)
            if heading_number_pattern.match(text) and i + 1 < len(doc.paragraphs) and not figure_line_pattern.match(doc.paragraphs[i + 1].text.strip()):
                i += 1
                next_text = doc.paragraphs[i].text.strip().replace('_', '.')
                run = new_doc.add_paragraph().add_run(next_text)
                run.bold = True
                run.font.size = Pt(12)
        else:
            text_content = [text]
            while i + 1 < len(doc.paragraphs) and not single_line_heading_pattern.match(doc.paragraphs[i + 1].text.strip()) and not heading_number_pattern.match(doc.paragraphs[i + 1].text.strip()) and not figure_line_pattern.match(doc.paragraphs[i + 1].text.strip()):
                i += 1
                next_text = doc.paragraphs[i].text.strip().replace('_', '.')
                text_content.append(next_text)
            consolidated_text = ' '.join(text_content)
            new_paragraph = new_doc.add_paragraph(consolidated_text)
            new_paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        i += 1

    new_doc.save(output_path)


def split_text_into_parts(text, limit):
    """
    Split text into parts, ensuring each part has around 'limit' words without cutting sentences in the middle.
    """
    sentences = re.split(r'(?<=[.!?])\s+', text)  # Split text into sentences
    parts = []
    part_words = []
    current_count = 0

    for sentence in sentences:
        sentence_words = sentence.split()
        sentence_length = len(sentence_words)

        if current_count + sentence_length > limit:
            # If adding this sentence exceeds the limit, start a new part
            if part_words:  # Only add to parts if there's content
                parts.append(' '.join(part_words))
            part_words = sentence_words  # Start a new part with the current sentence
            current_count = sentence_length
        else:
            # Otherwise, add this sentence to the current part
            part_words.extend(sentence_words)
            current_count += sentence_length

    # Add the last part if it has content
    if part_words:
        parts.append(' '.join(part_words))

    return parts


def customize_title_shape(title_shape, title_text):
    # Split the title text into sentences
    sentences = re.split(r'(?<=\.)\s+', title_text)
    # Iterate over each sentence and customize its appearance
    for index, sentence in enumerate(sentences):
        if sentence:
            p = title_shape.text_frame.add_paragraph() if index > 0 or title_shape.text_frame.paragraphs[0].text else title_shape.text_frame.paragraphs[0]
            p.text = sentence.replace('_', '.')
            # Calculate the font size based on the length of the text
            font_size = 28  # Default font size
            print(len(sentence), "#"*30)
            
            if 41 > len(sentence) > 31:  # Adjust as needed
                font_size = 24
            elif 51 >len(sentence) >= 41:
                font_size = 22
            elif len(sentence) >= 51:
                font_size = 18
            print(font_size, "@"*30)
            p.alignment = PP_ALIGN.JUSTIFY
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.name = 'Calibri'
    # Set shape options
    title_shape.width = Cm(21)
    title_shape.height = Cm(2.5)
    title_shape.left = Cm(0.3)
    title_shape.top = Cm(0.4)

def customize_content_shape(content_shape):
    content_shape.width = Cm(24)
    content_shape.height = Cm(15)
    content_shape.left = Cm(0.3)
    content_shape.top = Cm(2.8)


def add_content_to_slide(text_frame, text):
    sentences = re.split(r'(?<=\.)\s+', text)
    for index, sentence in enumerate(sentences):
        if sentence:
            p = text_frame.add_paragraph() if index > 0 or text_frame.paragraphs[0].text else text_frame.paragraphs[0]
            p.text = sentence.replace('_', '.')
            p.alignment = PP_ALIGN.JUSTIFY
            for run in p.runs:
                run.font.size = Pt(21)
                run.font.name = 'Calibri'


def add_slide_number(slide, number):
    # Define slide number position and size
    left = Cm(0.25)
    top = Cm(18)
    width = Cm(2)
    height = Cm(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(12)
    run.font.name = 'Calibri'


def process_pdf_to_word(pdf_path, image_folder_path, output_word_file):
    # Convert PDF pages to images
    image_paths = convert_pdf_pages_to_images(pdf_path, image_folder_path)
    
    # Detect highlighted regions in these images
    highlighted_regions = detect_highlighted_regions(image_paths)
    
    # Extract text from the highlighted regions
    extracted_texts = extract_text_from_highlights(highlighted_regions)
    
    # Save the extracted text to a Word document
    save_text_to_word(extracted_texts, output_word_file)


def run_process(pdf_path, image_folder_path, output_dir, template_path, start_button, status_label):
    
    # Disable the start button and update status
    start_button.config(state="disabled")
    status_label.config(text="Processing started...")
    
    def process():
        try:
            output_word_file = os.path.join(output_dir, "extracted_text.docx")
            formatted_output_word_file = os.path.join(output_dir, "formatted_extracted_text.docx")
            pptx_path = os.path.join(output_dir, "presentation.pptx")
            
            status_label.config(text="Converting PDF to Word...")
            process_pdf_to_word(pdf_path, image_folder_path, output_word_file)
            
            status_label.config(text="Formatting Word document...")
            format_headings_in_word(output_word_file, formatted_output_word_file)
            
            status_label.config(text="Creating PowerPoint presentation...")
            create_presentation_from_word(formatted_output_word_file, pptx_path, template_path)
            
            status_label.config(text="Processing completed successfully. Check the output directory.")
            messagebox.showinfo("Complete", "Processing completed successfully. Check the output directory.")
        except Exception as e:
            messagebox.showerror("Error", str(e))
            status_label.config(text="An error occurred. Please try again.")
        finally:
            # Re-enable the start button
            start_button.config(state="normal")
    
    # Run the processing in a separate thread to avoid freezing the GUI
    threading.Thread(target=process).start()


# GUI setup
root = tk.Tk()
root.title("PDF to PPT Converter")

tk.Label(root, text="Select PDF:").grid(row=0, column=0, padx=10, pady=10)
entry_pdf = tk.Entry(root, width=50)
entry_pdf.grid(row=0, column=1, padx=10, pady=10)
tk.Button(root, text="Browse", command=lambda: select_file(entry_pdf)).grid(row=0, column=2, padx=10, pady=10)

tk.Label(root, text="Image Folder:").grid(row=1, column=0, padx=10, pady=10)
entry_image_folder = tk.Entry(root, width=50)
entry_image_folder.grid(row=1, column=1, padx=10, pady=10)
tk.Button(root, text="Browse", command=lambda: select_folder("Select Image Folder", entry_image_folder)).grid(row=1, column=2, padx=10, pady=10)

tk.Label(root, text="Output Directory:").grid(row=2, column=0, padx=10, pady=10)
entry_output_dir = tk.Entry(root, width=50)
entry_output_dir.grid(row=2, column=1, padx=10, pady=10)
tk.Button(root, text="Browse", command=lambda: select_folder("Select Output Directory", entry_output_dir)).grid(row=2, column=2, padx=10, pady=10)

tk.Label(root, text="PPT Template:").grid(row=3, column=0, padx=10, pady=10)
entry_template = tk.Entry(root, width=50)
entry_template.grid(row=3, column=1, padx=10, pady=10)
tk.Button(root, text="Browse", command=lambda: select_file(entry_template)).grid(row=3, column=2
                                                                                 , padx=10, pady=10)
status_label = tk.Label(root, text="Ready")
status_label.grid(row=5, column=1, padx=10, pady=10)
button_process = tk.Button(root, text="Start Processing", command=start_process_thread)
button_process.grid(row=4, column=1, padx=10, pady=10)

root.mainloop()



